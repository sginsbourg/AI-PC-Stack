# replace_speaker_notes.py
from pptx import Presentation
import re
from config_loader import config

def parse_speaker_notes():
    """
    Parse the Markdown file to extract speaker notes for each slide.
    Returns a dictionary with slide number as key and notes text as value.
    """
    md_file_path = config.get('INPUT_FILE')
    notes_dict = {}
    current_slide = None
    current_notes = []
    
    with open(md_file_path, 'r', encoding='utf-8') as file:
        for line in file:
            line = line.strip()
            match = re.match(r'##\s*Slide\s*(\d+):', line)
            if match:
                if current_slide is not None and current_notes:
                    notes_dict[current_slide] = '\n'.join(current_notes).strip()
                current_slide = int(match.group(1))
                current_notes = []
            elif current_slide is not None:
                current_notes.append(line)
    
    if current_slide is not None and current_notes:
        notes_dict[current_slide] = '\n'.join(current_notes).strip()
    
    return notes_dict

def replace_speaker_notes():
    """
    Open the PowerPoint file, replace speaker notes for each slide, and save to output file.
    """
    pptx_file_path = config.get('PPTX_FILE')
    output_pptx_path = config.get('UPDATED_PPTX')
    notes_dict = parse_speaker_notes()
    
    prs = Presentation(pptx_file_path)
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        if slide_num in notes_dict:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_dict[slide_num]
    
    prs.save(output_pptx_path)
    print(f"Speaker notes replaced and saved to {output_pptx_path}")

if __name__ == "__main__":
    replace_speaker_notes()